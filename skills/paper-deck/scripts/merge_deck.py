#!/usr/bin/env python3
"""
Merge generated slide images into PPTX and PDF.
Supports:
  - Mixed text overlay on AIGC backgrounds
  - Generation log verification
  - Image resolution / aspect-ratio validation
  - Outline-vs-image count consistency check

Usage:
    python3 merge_deck.py /path/to/paper-deck/topic-slug
    python3 merge_deck.py /path/to/paper-deck/topic-slug --check
    python3 merge_deck.py /path/to/paper-deck/topic-slug --overlay overlays.json
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# ── Constants ──────────────────────────────────────────────────────────────

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
SLIDE_W_IN = 13.333333  # 16:9
SLIDE_H_IN = 7.5
SLIDE_W_PX = 1280
SLIDE_H_PX = 720
MIN_IMAGE_W_PX = 800
MAX_RATIO_DEVIATION = 0.05  # ±5%


# ── Data Models ─────────────────────────────────────────────────────────────

@dataclass
class GenerationRecord:
    slide_number: int
    prompt_file: str
    output_file: str
    backend: str = "unknown"
    timestamp: str = ""


@dataclass
class TextOverlay:
    slide_number: int
    text: str
    x: float  # fraction 0-1 of slide width
    y: float  # fraction 0-1 of slide height
    font_size: int = 32
    color: str = "#FFFFFF"
    align: str = "left"


@dataclass
class QualityReport:
    passed: bool = True
    checks: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)


# ── Helpers ─────────────────────────────────────────────────────────────────

def natural_key(path: Path) -> tuple[int, str]:
    match = re.match(r"^(\d+)", path.stem)
    number = int(match.group(1)) if match else 9999
    return number, path.name


def find_images(deck_dir: Path) -> list[Path]:
    image_dir = deck_dir / "images"
    if not image_dir.exists():
        raise SystemExit(f"[ERROR] Missing images directory: {image_dir}")

    images = [p for p in image_dir.iterdir() if p.suffix.lower() in IMAGE_EXTS]
    images.sort(key=natural_key)
    if not images:
        raise SystemExit(f"[ERROR] No slide images found in: {image_dir}")
    return images


def load_overlays(path: Optional[Path]) -> dict[int, list[TextOverlay]]:
    if not path or not path.exists():
        return {}
    with open(path, encoding="utf-8") as f:
        data = json.load(f)
    overlays: dict[int, list[TextOverlay]] = {}
    for item in data.get("overlays", []):
        sn = item["slide_number"]
        overlays.setdefault(sn, []).append(TextOverlay(
            slide_number=sn,
            text=item["text"],
            x=item.get("x", 0.5),
            y=item.get("y", 0.5),
            font_size=item.get("font_size", 32),
            color=item.get("color", "#FFFFFF"),
            align=item.get("align", "left"),
        ))
    return overlays


def load_generation_log(deck_dir: Path) -> list[GenerationRecord]:
    log_path = deck_dir / "generation-log.md"
    if not log_path.exists():
        return []
    records: list[GenerationRecord] = []
    with open(log_path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            parts = [p.strip() for p in line.strip("|").split("|")]
            if len(parts) >= 4 and parts[0].isdigit():
                records.append(GenerationRecord(
                    slide_number=int(parts[0]),
                    prompt_file=parts[1],
                    output_file=parts[2],
                    backend=parts[3] if len(parts) > 3 else "unknown",
                    timestamp=parts[4] if len(parts) > 4 else "",
                ))
    return records


# ── Quality Checks ──────────────────────────────────────────────────────────

def check_images(images: list[Path]) -> QualityReport:
    report = QualityReport()
    for img_path in images:
        try:
            with Image.open(img_path) as img:
                w, h = img.size
                if w < MIN_IMAGE_W_PX:
                    report.warnings.append(
                        f"{img_path.name}: {w}x{h} — below {MIN_IMAGE_W_PX}px wide"
                    )
                actual_ratio = w / h
                expected_ratio = SLIDE_W_PX / SLIDE_H_PX
                if abs(actual_ratio - expected_ratio) > MAX_RATIO_DEVIATION:
                    report.warnings.append(
                        f"{img_path.name}: aspect ratio {w/h:.2f} deviates from 16:9"
                    )
                report.checks.append(f"{img_path.name}: {w}x{h} OK")
        except Exception as e:
            report.errors.append(f"{img_path.name}: cannot open — {e}")
            report.passed = False
    return report


def check_generation_log(images: list[Path], deck_dir: Path) -> QualityReport:
    report = QualityReport()
    records = load_generation_log(deck_dir)
    img_map = {p.name for p in images}
    logged_files = {r.output_file for r in records}
    for img_name in img_map:
        if img_name not in logged_files:
            report.warnings.append(f"{img_name}: no generation-log.md entry")
    report.checks.append(f"Generation log records: {len(records)}")
    return report


def check_outline_vs_images(images: list[Path], deck_dir: Path) -> QualityReport:
    report = QualityReport()
    outline_path = deck_dir / "outline.md"
    if not outline_path.exists():
        report.warnings.append("No outline.md found")
        return report
    with open(outline_path, encoding="utf-8") as f:
        content = f.read()
    slide_headers = re.findall(r"^##\s+\d+\.", content, re.MULTILINE)
    if len(slide_headers) != len(images):
        report.warnings.append(
            f"outline.md: {len(slide_headers)} slides, images/: {len(images)}"
        )
    else:
        report.checks.append(f"Slide count: {len(images)}")
    return report


# ── PPTX / PDF Builders ────────────────────────────────────────────────────

def make_pptx(
    images: list[Path],
    output: Path,
    overlays: dict[int, list[TextOverlay]],
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    for idx, image in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image), 0, 0,
            width=prs.slide_width, height=prs.slide_height,
        )
        slide_num = idx + 1
        for overlay in overlays.get(slide_num, []):
            txBox = slide.shapes.add_textbox(
                Inches(overlay.x * SLIDE_W_IN),
                Inches(overlay.y * SLIDE_H_IN),
                Inches(SLIDE_W_IN * 0.8),
                Inches(1.0),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = overlay.text
            p.font.size = Pt(overlay.font_size)
            p.font.color.rgb = _parse_color(overlay.color)
            from pptx.enum.text import PP_ALIGN
            p.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(overlay.align, PP_ALIGN.LEFT)

    prs.save(output)


def make_pdf(images: list[Path], output: Path) -> None:
    frames: list[Image.Image] = []
    for image_path in images:
        with Image.open(image_path) as img:
            frames.append(img.convert("RGB").copy())
    first, *rest = frames
    first.save(output, save_all=True, append_images=rest, resolution=150.0)


def _parse_color(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0xFF, 0xFF, 0xFF)
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


# ── CLI ─────────────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Merge slide images into PPTX/PDF with overlays + quality checks."
    )
    parser.add_argument("deck_dir", help="Deck directory containing images/")
    parser.add_argument("--name", help="Output base name (default: directory name)")
    parser.add_argument("--check", action="store_true",
                        help="Run quality checks without merging")
    parser.add_argument("--overlay", type=Path, default=None,
                        help="JSON file with text overlay definitions")
    parser.add_argument("--skip-checks", action="store_true",
                        help="Skip quality checks before merge")
    args = parser.parse_args()

    deck_dir = Path(args.deck_dir).expanduser().resolve()
    if not deck_dir.exists():
        raise SystemExit(f"[ERROR] Deck directory does not exist: {deck_dir}")

    images = find_images(deck_dir)
    overlays = load_overlays(args.overlay)

    # ── Check-only mode ──
    if args.check:
        print(f"\n{'='*50}")
        print(f"Quality Check — {deck_dir.name}")
        print(f"{'='*50}\n")
        passed = True
        for fn in [check_images, check_generation_log, check_outline_vs_images]:
            report = fn(images, deck_dir)
            for c in report.checks:
                print(f"  [OK] {c}")
            for w in report.warnings:
                print(f"  [WARN] {w}")
            for e in report.errors:
                print(f"  [FAIL] {e}")
                passed = False
        print(f"\nResult: {'PASS' if passed else 'ISSUES FOUND'}")
        sys.exit(0 if passed else 1)

    # ── Pre-merge quality gates ──
    if not args.skip_checks:
        all_pass = True
        for fn in [check_images, check_generation_log, check_outline_vs_images]:
            report = fn(images, deck_dir)
            for w in report.warnings:
                print(f"  [WARN] {w}")
            for e in report.errors:
                print(f"  [FAIL] {e}")
                all_pass = False
        if not all_pass:
            print("[ERROR] Quality checks failed. Use --skip-checks to force.")
            sys.exit(1)

    # ── Merge ──
    base = args.name or deck_dir.name
    pptx_path = deck_dir / f"{base}.pptx"
    pdf_path = deck_dir / f"{base}.pdf"

    make_pptx(images, pptx_path, overlays)
    make_pdf(images, pdf_path)

    print(f"\n{'='*50}")
    print(f"Merged {len(images)} slides")
    if overlays:
        print(f"Overlays: {sum(len(v) for v in overlays.values())}")
    print(f"PPTX: {pptx_path}")
    print(f"PDF:  {pdf_path}")
    print(f"{'='*50}")


if __name__ == "__main__":
    main()
